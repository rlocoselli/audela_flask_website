from __future__ import annotations

import io
from decimal import Decimal
from typing import Any
from datetime import datetime
from .export_branding import resolve_brand_tokens


def _to_number(value: Any) -> float | None:
    if value is None or value == "":
        return None
    if isinstance(value, (int, float, Decimal)):
        try:
            return float(value)
        except Exception:
            return None
    txt = str(value).strip()
    if not txt:
        return None
    txt = txt.replace(" ", "")
    if "," in txt and "." in txt:
        if txt.rfind(",") > txt.rfind("."):
            txt = txt.replace(".", "").replace(",", ".")
        else:
            txt = txt.replace(",", "")
    elif "," in txt:
        txt = txt.replace(",", ".")
    try:
        return float(txt)
    except Exception:
        return None


def _style_tokens(style_guide: str) -> dict[str, Any]:
    b = resolve_brand_tokens(style_guide)
    return {
        "bg": str(b.get("bg") or "#F5F7FB").lstrip("#"),
        "accent": str(b.get("accent") or "#1E5CC6").lstrip("#"),
        "accent2": str(b.get("accent2") or "#0FA3B1").lstrip("#"),
        "ink": str(b.get("ink") or "#0E1A2B").lstrip("#"),
        "muted": str(b.get("muted") or "#4F5E73").lstrip("#"),
        "font": str(b.get("font_ppt") or "Calibri"),
    }


def _metric_cards(columns: list[str], rows: list[list[Any]], max_items: int = 4) -> list[dict[str, Any]]:
    if not columns or not rows:
        return []
    out: list[dict[str, Any]] = []
    for idx, col in enumerate(columns[:12]):
        vals: list[float] = []
        for r in rows[:500]:
            if not isinstance(r, (list, tuple)) or idx >= len(r):
                continue
            n = _to_number(r[idx])
            if n is not None:
                vals.append(n)
        if len(vals) < 3:
            continue
        out.append(
            {
                "label": str(col),
                "avg": sum(vals) / max(1, len(vals)),
                "min": min(vals),
                "max": max(vals),
            }
        )
        if len(out) >= max_items:
            break
    return out


def _top_dimension_breakdown(columns: list[str], rows: list[list[Any]], max_items: int = 8) -> tuple[str, str, list[tuple[str, float]]]:
    if not columns or not rows:
        return "", "", []

    numeric_idx = -1
    for idx, _ in enumerate(columns):
        vals = []
        for r in rows[:400]:
            if not isinstance(r, (list, tuple)) or idx >= len(r):
                continue
            n = _to_number(r[idx])
            if n is not None:
                vals.append(n)
        if len(vals) >= 5:
            numeric_idx = idx
            break

    if numeric_idx < 0:
        return "", "", []

    dim_idx = -1
    for idx, _ in enumerate(columns):
        if idx == numeric_idx:
            continue
        non_num = 0
        seen = 0
        for r in rows[:300]:
            if not isinstance(r, (list, tuple)) or idx >= len(r):
                continue
            v = r[idx]
            if v in (None, ""):
                continue
            seen += 1
            if _to_number(v) is None:
                non_num += 1
        if seen and (non_num / seen) >= 0.6:
            dim_idx = idx
            break

    if dim_idx < 0:
        return "", "", []

    agg: dict[str, float] = {}
    for r in rows[:3000]:
        if not isinstance(r, (list, tuple)) or dim_idx >= len(r) or numeric_idx >= len(r):
            continue
        key = str(r[dim_idx] if r[dim_idx] not in (None, "") else "(empty)").strip()[:42]
        n = _to_number(r[numeric_idx])
        if n is None:
            continue
        agg[key] = float(agg.get(key, 0.0) + float(n))

    top = sorted(agg.items(), key=lambda kv: kv[1], reverse=True)[: max(1, int(max_items or 8))]
    return str(columns[dim_idx]), str(columns[numeric_idx]), top


def table_to_pptx_bytes(
    *,
    title: str,
    source_name: str,
    analysis: str,
    columns: list[str],
    rows: list[list[Any]],
    style_guide: str = "",
) -> bytes:
    """Build a simple, styled PowerPoint deck from tabular BI output.

    Slides:
    - Cover
    - Executive summary
    - KPI cards
    - Top category breakdown
    - Data sample table
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.util import Inches, Pt
    except Exception as e:
        raise RuntimeError("python-pptx is not installed. Please add python-pptx to requirements.") from e

    prs = Presentation()
    theme = _style_tokens(style_guide)
    generated_at = datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")

    def rgb(hex_color: str):
        return RGBColor.from_string(str(hex_color or "000000"))

    def apply_chrome(slide, *, section_title: str, slide_no: int):
        top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.33), Inches(0.18))
        top.fill.solid()
        top.fill.fore_color.rgb = rgb(theme["accent"])
        top.line.fill.background()

        footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.25)).text_frame
        footer.clear()
        p = footer.paragraphs[0]
        p.text = f"{section_title}  ·  Generated {generated_at}  ·  Slide {slide_no}"
        p.font.size = Pt(9)
        p.font.color.rgb = rgb(theme["muted"])
        p.font.name = theme["font"]

    # Slide 1: Cover
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    bg = cover.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(theme["bg"])

    banner = cover.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.35))
    banner.fill.solid()
    banner.fill.fore_color.rgb = rgb(theme["accent"])
    banner.line.fill.background()
    tf = banner.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(title or "Executive Deck")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = theme["font"]

    sub = cover.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11), Inches(1.0)).text_frame
    sub.clear()
    p2 = sub.paragraphs[0]
    p2.text = f"Source: {source_name or 'BI Lite data'}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = rgb(theme["ink"])
    p2.font.name = theme["font"]

    chip = cover.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.05), Inches(5.6), Inches(0.55))
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(theme["accent2"])
    chip.line.fill.background()
    chip_tf = chip.text_frame
    chip_tf.clear()
    cp = chip_tf.paragraphs[0]
    cp.text = f"Rows: {len(rows):,}  |  Columns: {len(columns):,}"
    cp.font.size = Pt(12)
    cp.font.bold = True
    cp.font.color.rgb = RGBColor(255, 255, 255)
    cp.font.name = theme["font"]
    apply_chrome(cover, section_title="Cover", slide_no=1)

    # Slide 2: Executive summary
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    b2 = s2.background.fill
    b2.solid()
    b2.fore_color.rgb = rgb(theme["bg"])

    h2 = s2.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(0.7)).text_frame
    h2.clear()
    p = h2.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = rgb(theme["accent"])
    p.font.name = theme["font"]

    body = s2.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.9)).text_frame
    body.word_wrap = True
    body.clear()
    analysis_text = str(analysis or "No AI summary available. Use the data table and KPIs on next slides.")
    summary_lines = [x.strip() for x in analysis_text.splitlines() if x.strip()] or [analysis_text]
    for idx, line in enumerate(summary_lines[:7]):
        para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        para.text = ("• " if idx > 0 else "") + line[:360]
        para.level = 0
        para.font.size = Pt(15 if idx else 16)
        para.font.color.rgb = rgb(theme["ink"])
        para.font.name = theme["font"]
    apply_chrome(s2, section_title="Executive Summary", slide_no=2)

    # Slide 3: KPI cards
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    b3 = s3.background.fill
    b3.solid()
    b3.fore_color.rgb = rgb(theme["bg"])

    h3 = s3.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(0.7)).text_frame
    h3.clear()
    p = h3.paragraphs[0]
    p.text = "Key Metrics"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = rgb(theme["accent"])
    p.font.name = theme["font"]

    metrics = _metric_cards(columns, rows, max_items=4)
    if not metrics:
        note = s3.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(10.5), Inches(1.0)).text_frame
        note.clear()
        p = note.paragraphs[0]
        p.text = "No numeric columns detected for KPI cards."
        p.font.size = Pt(16)
        p.font.color.rgb = rgb(theme["muted"])
        p.font.name = theme["font"]
    else:
        for i, m in enumerate(metrics[:4]):
            col = i % 2
            row = i // 2
            x = Inches(0.8 + (col * 6.2))
            y = Inches(1.4 + (row * 2.5))
            card = s3.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.1))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = rgb(theme["accent2"])

            tf = card.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = str(m.get("label") or "Metric")[:50]
            p0.font.size = Pt(14)
            p0.font.bold = True
            p0.font.color.rgb = rgb(theme["ink"])
            p0.font.name = theme["font"]

            p1 = tf.add_paragraph()
            p1.text = f"Avg: {float(m.get('avg') or 0.0):,.2f}"
            p1.font.size = Pt(20)
            p1.font.bold = True
            p1.font.color.rgb = rgb(theme["accent"])
            p1.font.name = theme["font"]

            p2 = tf.add_paragraph()
            p2.text = f"Min: {float(m.get('min') or 0.0):,.2f}   Max: {float(m.get('max') or 0.0):,.2f}"
            p2.font.size = Pt(12)
            p2.font.color.rgb = rgb(theme["muted"])
            p2.font.name = theme["font"]
    apply_chrome(s3, section_title="Key Metrics", slide_no=3)

    # Slide 4: Top category breakdown
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    b4 = s4.background.fill
    b4.solid()
    b4.fore_color.rgb = rgb(theme["bg"])

    h4 = s4.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.7)).text_frame
    h4.clear()
    p = h4.paragraphs[0]
    p.text = "Top Category Breakdown"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = rgb(theme["accent"])
    p.font.name = theme["font"]

    dim_name, metric_name, top_items = _top_dimension_breakdown(columns, rows, max_items=8)
    if top_items:
        max_val = max(float(v) for _, v in top_items) if top_items else 0.0
        y = 1.4
        subtitle = s4.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(11), Inches(0.35)).text_frame
        subtitle.clear()
        sp = subtitle.paragraphs[0]
        sp.text = f"{dim_name} by {metric_name}"
        sp.font.size = Pt(12)
        sp.font.color.rgb = rgb(theme["muted"])
        sp.font.name = theme["font"]

        for label, val in top_items[:8]:
            lbox = s4.shapes.add_textbox(Inches(0.8), Inches(y), Inches(3.1), Inches(0.3)).text_frame
            lbox.clear()
            lp = lbox.paragraphs[0]
            lp.text = str(label)
            lp.font.size = Pt(11)
            lp.font.color.rgb = rgb(theme["ink"])
            lp.font.name = theme["font"]

            w = 0.1 if max_val <= 0 else max(0.3, min(6.8, (float(val) / max_val) * 6.8))
            bar = s4.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.9), Inches(y + 0.01), Inches(w), Inches(0.23))
            bar.fill.solid()
            bar.fill.fore_color.rgb = rgb(theme["accent2"])
            bar.line.fill.background()

            vbox = s4.shapes.add_textbox(Inches(10.9), Inches(y), Inches(1.7), Inches(0.3)).text_frame
            vbox.clear()
            vp = vbox.paragraphs[0]
            vp.text = f"{float(val):,.2f}"
            vp.font.size = Pt(10)
            vp.font.bold = True
            vp.font.color.rgb = rgb(theme["ink"])
            vp.font.name = theme["font"]
            y += 0.45
    else:
        note = s4.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(10.8), Inches(1.0)).text_frame
        note.clear()
        p = note.paragraphs[0]
        p.text = "Could not infer a categorical breakdown from the current dataset."
        p.font.size = Pt(14)
        p.font.color.rgb = rgb(theme["muted"])
        p.font.name = theme["font"]
    apply_chrome(s4, section_title="Category Breakdown", slide_no=4)

    # Slide 5: Data sample
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    b5 = s5.background.fill
    b5.solid()
    b5.fore_color.rgb = rgb(theme["bg"])

    h5 = s5.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.7)).text_frame
    h5.clear()
    p = h5.paragraphs[0]
    p.text = "Data Snapshot"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = rgb(theme["accent"])
    p.font.name = theme["font"]

    max_cols = min(6, len(columns))
    max_rows = min(12, len(rows))
    if max_cols > 0:
        table_shape = s5.shapes.add_table(max_rows + 1, max_cols, Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.6))
        table = table_shape.table

        for c in range(max_cols):
            cell = table.cell(0, c)
            cell.text = str(columns[c])[:28]
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(theme["accent"])
            for para in cell.text_frame.paragraphs:
                para.font.bold = True
                para.font.size = Pt(11)
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.name = theme["font"]

        for r in range(max_rows):
            row_data = rows[r] if isinstance(rows[r], (list, tuple)) else []
            for c in range(max_cols):
                val = row_data[c] if c < len(row_data) else ""
                cell = table.cell(r + 1, c)
                cell.text = str(val)[:40]
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(10)
                    para.font.color.rgb = rgb(theme["ink"])
                    para.font.name = theme["font"]
                if r % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 248, 252)
    apply_chrome(s5, section_title="Data Snapshot", slide_no=5)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
